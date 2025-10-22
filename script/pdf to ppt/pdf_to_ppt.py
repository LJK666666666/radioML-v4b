#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PDF to PowerPoint Converter
将PDF文件转换为PowerPoint演示文稿，每页PDF作为一张图片插入到PPT中
"""

import os
import sys
import time
from pathlib import Path
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import tempfile

def pdf_to_ppt(pdf_path, output_path=None, dpi=200):
    """
    将PDF文件转换为PowerPoint演示文稿
    
    Args:
        pdf_path (str): PDF文件路径
        output_path (str): 输出PPT文件路径，如果为None则使用PDF文件名
        dpi (int): 图片分辨率，默认200
    """
    
    # 检查PDF文件是否存在
    if not os.path.exists(pdf_path):
        raise FileNotFoundError(f"PDF文件不存在: {pdf_path}")
    
    # 设置输出文件路径
    if output_path is None:
        pdf_name = Path(pdf_path).stem
        output_path = f"{pdf_name}.pptx"
    
    print(f"开始转换PDF文件: {pdf_path}")
    print(f"输出文件: {output_path}")
    
    try:
        # 将PDF转换为图片列表
        print("正在将PDF转换为图片...")
        images = convert_from_path(pdf_path, dpi=dpi)
        print(f"成功转换 {len(images)} 页")
        
        # 创建PowerPoint演示文稿
        prs = Presentation()
        
        # 获取幻灯片布局（空白布局）
        blank_slide_layout = prs.slide_layouts[6]  # 空白布局
        
        # 为每个图片创建一张幻灯片
        for i, image in enumerate(images):
            print(f"正在处理第 {i+1} 页...")
            
            # 添加幻灯片
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # 创建临时文件保存图片
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
            temp_file_path = temp_file.name
            temp_file.close()  # 先关闭文件句柄
            
            try:
                # 保存图片到临时文件
                image.save(temp_file_path, 'PNG')
                
                # 获取图片尺寸（像素）
                img_width, img_height = image.size
                print(f"  PDF页面尺寸: {img_width} x {img_height} 像素")
                
                # 根据PDF页面的实际DPI计算真实尺寸（英寸）
                # 使用用户指定的DPI作为转换基准
                img_width_inch = img_width / dpi
                img_height_inch = img_height / dpi
                print(f"  实际页面尺寸: {img_width_inch:.2f} x {img_height_inch:.2f} 英寸")
                
                # 设置幻灯片尺寸为PDF页面的实际尺寸
                slide_width_inch = img_width_inch
                slide_height_inch = img_height_inch
                
                # 转换为PPT内部单位
                slide_width = Inches(slide_width_inch)
                slide_height = Inches(slide_height_inch)
                
                # 更新演示文稿的幻灯片尺寸
                prs.slide_width = slide_width
                prs.slide_height = slide_height
                print(f"  设置幻灯片尺寸: {slide_width_inch:.2f} x {slide_height_inch:.2f} 英寸")
                
                # 图片占满整个幻灯片
                final_width = slide_width
                final_height = slide_height
                left = 0
                top = 0
                
                # 添加图片到幻灯片
                slide.shapes.add_picture(
                    temp_file_path, 
                    left, 
                    top, 
                    final_width, 
                    final_height
                )
                
            finally:
                # 确保删除临时文件，添加重试机制
                for retry in range(3):
                    try:
                        if os.path.exists(temp_file_path):
                            os.unlink(temp_file_path)
                        break
                    except Exception as e:
                        if retry < 2:  # 前两次失败时等待后重试
                            time.sleep(0.1)
                            continue
                        else:
                            print(f"警告: 无法删除临时文件 {temp_file_path}: {e}")
                            # 不抛出异常，继续处理
        
        # 保存PowerPoint文件
        prs.save(output_path)
        print(f"转换完成！输出文件: {output_path}")
        print(f"共创建了 {len(images)} 张幻灯片")
        
    except Exception as e:
        print(f"转换过程中出现错误: {str(e)}")
        raise

def main():
    """主函数"""
    # 获取当前脚本所在目录
    script_dir = os.path.dirname(os.path.abspath(__file__))
    
    # 默认PDF文件路径
    default_pdf_path = os.path.join(script_dir, "presentation.pdf")
    
    # 检查命令行参数
    if len(sys.argv) > 1:
        pdf_path = sys.argv[1]
    else:
        pdf_path = default_pdf_path
    
    # 检查PDF文件是否存在
    if not os.path.exists(pdf_path):
        print(f"错误: PDF文件不存在: {pdf_path}")
        print("用法: python pdf_to_ppt.py [PDF文件路径]")
        print(f"或者将PDF文件命名为 'presentation.pdf' 并放在脚本同目录下")
        sys.exit(1)
    
    try:
        # 执行转换
        pdf_to_ppt(pdf_path)
        print("\n转换成功完成！")
        
    except Exception as e:
        print(f"\n转换失败: {str(e)}")
        sys.exit(1)

if __name__ == "__main__":
    main()
