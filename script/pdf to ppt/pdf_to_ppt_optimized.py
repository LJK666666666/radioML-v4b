#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PDF to PowerPoint Converter - 优化版本
修复了临时文件访问问题，提高了稳定性和性能
"""

import os
import sys
import time
import gc
from pathlib import Path
from pdf2image import convert_from_path
from pptx import Presentation
from PIL import Image
import tempfile
import io

def pdf_to_ppt_optimized(pdf_path, output_path=None, dpi=200, batch_size=5):
    """
    优化的PDF转PowerPoint转换器
    
    Args:
        pdf_path (str): PDF文件路径
        output_path (str): 输出PPT文件路径，如果为None则使用PDF文件名
        dpi (int): 图片分辨率，默认200
        batch_size (int): 批处理大小，避免内存问题
    """
    
    # 检查PDF文件是否存在
    if not os.path.exists(pdf_path):
        raise FileNotFoundError(f"PDF文件不存在: {pdf_path}")
    
    # 设置输出文件路径
    if output_path is None:
        pdf_name = Path(pdf_path).stem
        output_path = f"{pdf_name}_optimized.pptx"
    
    print(f"开始转换PDF文件: {pdf_path}")
    print(f"输出文件: {output_path}")
    print(f"分辨率: {dpi} DPI")
    
    try:
        # 获取PDF页数（不转换全部，只获取页数信息）
        print("正在分析PDF文件...")
        sample_images = convert_from_path(pdf_path, dpi=72, first_page=1, last_page=1)
        if not sample_images:
            raise ValueError("无法读取PDF文件或文件为空")
        
        # 获取总页数
        total_pages = len(convert_from_path(pdf_path, dpi=72))
        print(f"PDF总页数: {total_pages}")
        
        # 创建PowerPoint演示文稿
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]  # 空白布局
        
        # 分批处理页面以避免内存问题
        for start_page in range(1, total_pages + 1, batch_size):
            end_page = min(start_page + batch_size - 1, total_pages)
            
            print(f"正在处理第 {start_page}-{end_page} 页...")
            
            # 转换当前批次的页面
            batch_images = convert_from_path(
                pdf_path, 
                dpi=dpi, 
                first_page=start_page, 
                last_page=end_page
            )
            
            # 处理批次中的每一页
            for page_index, image in enumerate(batch_images):
                current_page = start_page + page_index
                print(f"  处理第 {current_page} 页...")
                
                # 添加幻灯片
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # 使用内存流而不是临时文件
                img_stream = io.BytesIO()
                image.save(img_stream, format='PNG', optimize=True)
                img_stream.seek(0)
                
                # 创建临时文件（更安全的方式）
                temp_fd, temp_path = tempfile.mkstemp(suffix='.png')
                try:
                    # 写入临时文件
                    with os.fdopen(temp_fd, 'wb') as temp_file:
                        temp_file.write(img_stream.getvalue())
                    
                    # 获取图片尺寸
                    img_width, img_height = image.size
                    
                    # 计算幻灯片尺寸
                    slide_width = prs.slide_width
                    slide_height = prs.slide_height
                    
                    # 计算最佳缩放比例
                    width_ratio = slide_width / img_width
                    height_ratio = slide_height / img_height
                    scale_ratio = min(width_ratio, height_ratio)
                    
                    # 计算最终尺寸和位置
                    final_width = int(img_width * scale_ratio)
                    final_height = int(img_height * scale_ratio)
                    left = (slide_width - final_width) // 2
                    top = (slide_height - final_height) // 2
                    
                    # 添加图片到幻灯片
                    slide.shapes.add_picture(
                        temp_path,
                        left,
                        top,
                        final_width,
                        final_height
                    )
                    
                finally:
                    # 安全删除临时文件
                    try:
                        os.unlink(temp_path)
                    except:
                        pass
                
                # 清理内存
                img_stream.close()
                del image
            
            # 清理批次数据
            del batch_images
            gc.collect()  # 强制垃圾回收
        
        # 保存PowerPoint文件
        print("正在保存PowerPoint文件...")
        prs.save(output_path)
        print(f"转换完成！输出文件: {output_path}")
        print(f"共创建了 {total_pages} 张幻灯片")
        
        return output_path
        
    except Exception as e:
        print(f"转换过程中出现错误: {str(e)}")
        raise

def main():
    """主函数"""
    # 获取当前脚本所在目录
    script_dir = os.path.dirname(os.path.abspath(__file__))
    
    # 默认PDF文件路径
    default_pdf_path = os.path.join(script_dir, "presentation.pdf")
    
    # 检查命令行参数
    if len(sys.argv) > 1:
        pdf_path = sys.argv[1]
        dpi = int(sys.argv[2]) if len(sys.argv) > 2 else 200
    else:
        pdf_path = default_pdf_path
        dpi = 200
    
    # 检查PDF文件是否存在
    if not os.path.exists(pdf_path):
        print(f"错误: PDF文件不存在: {pdf_path}")
        print("用法: python pdf_to_ppt_optimized.py [PDF文件路径] [DPI(可选)]")
        print(f"或者将PDF文件命名为 'presentation.pdf' 并放在脚本同目录下")
        sys.exit(1)
    
    try:
        start_time = time.time()
        
        # 执行转换
        result = pdf_to_ppt_optimized(pdf_path, dpi=dpi)
        
        end_time = time.time()
        elapsed_time = end_time - start_time
        
        print(f"\n转换成功完成！")
        print(f"耗时: {elapsed_time:.2f} 秒")
        print(f"输出文件: {result}")
        
    except Exception as e:
        print(f"\n转换失败: {str(e)}")
        print("\n故障排除提示:")
        print("1. 确保PDF文件未被其他程序占用")
        print("2. 确保有足够的磁盘空间")
        print("3. 尝试降低DPI设置 (例如: python script.py file.pdf 150)")
        print("4. 确保已正确安装所有依赖包")
        sys.exit(1)

if __name__ == "__main__":
    main()
