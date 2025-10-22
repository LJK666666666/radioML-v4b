#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
高级PDF to PowerPoint转换器
支持批量转换、进度显示、错误恢复等功能
"""

import os
import sys
import argparse
import json
from pathlib import Path
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import tempfile
from typing import List, Optional, Tuple
import time

class PDFToPPTConverter:
    """PDF to PowerPoint转换器类"""
    
    def __init__(self, dpi: int = 200, poppler_path: Optional[str] = None):
        """
        初始化转换器
        
        Args:
            dpi (int): 图片分辨率
            poppler_path (str): poppler工具路径（Windows需要）
        """
        self.dpi = dpi
        self.poppler_path = poppler_path
        
    def convert_single_pdf(self, pdf_path: str, output_path: Optional[str] = None, 
                          progress_callback=None) -> str:
        """
        转换单个PDF文件
        
        Args:
            pdf_path (str): PDF文件路径
            output_path (str): 输出PPT文件路径
            progress_callback: 进度回调函数
            
        Returns:
            str: 输出文件路径
        """
        
        # 检查PDF文件
        if not os.path.exists(pdf_path):
            raise FileNotFoundError(f"PDF文件不存在: {pdf_path}")
        
        # 设置输出路径
        if output_path is None:
            pdf_name = Path(pdf_path).stem
            output_dir = Path(pdf_path).parent
            output_path = output_dir / f"{pdf_name}.pptx"
        
        if progress_callback:
            progress_callback(f"开始转换: {Path(pdf_path).name}")
        
        try:
            # 转换PDF为图片
            if progress_callback:
                progress_callback("正在将PDF转换为图片...")
            
            convert_kwargs = {
                'pdf_path': pdf_path,
                'dpi': self.dpi,
                'fmt': 'RGB'
            }
            
            # 如果指定了poppler路径，添加到参数中
            if self.poppler_path:
                convert_kwargs['poppler_path'] = self.poppler_path
            
            images = convert_from_path(**convert_kwargs)
            
            if progress_callback:
                progress_callback(f"成功转换 {len(images)} 页")
            
            # 创建PowerPoint
            prs = Presentation()
            blank_slide_layout = prs.slide_layouts[6]  # 空白布局
            
            # 处理每一页
            for i, image in enumerate(images):
                if progress_callback:
                    progress_callback(f"处理第 {i+1}/{len(images)} 页")
                
                # 添加幻灯片
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # 保存图片到临时文件
                with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as temp_file:
                    image.save(temp_file.name, 'PNG', optimize=True)
                    
                    # 计算图片位置和尺寸
                    img_width, img_height = image.size
                    slide_width = prs.slide_width
                    slide_height = prs.slide_height
                    
                    # 保持宽高比的缩放
                    width_ratio = slide_width / img_width
                    height_ratio = slide_height / img_height
                    scale_ratio = min(width_ratio, height_ratio)
                    
                    final_width = int(img_width * scale_ratio)
                    final_height = int(img_height * scale_ratio)
                    
                    # 居中位置
                    left = (slide_width - final_width) // 2
                    top = (slide_height - final_height) // 2
                    
                    # 添加图片
                    slide.shapes.add_picture(
                        temp_file.name, left, top, final_width, final_height
                    )
                    
                    # 清理临时文件
                    os.unlink(temp_file.name)
            
            # 保存PPT
            if progress_callback:
                progress_callback("正在保存PowerPoint文件...")
            
            prs.save(str(output_path))
            
            if progress_callback:
                progress_callback(f"转换完成: {output_path}")
            
            return str(output_path)
            
        except Exception as e:
            raise RuntimeError(f"转换过程中出现错误: {str(e)}")
    
    def convert_batch(self, pdf_files: List[str], output_dir: Optional[str] = None) -> List[str]:
        """
        批量转换PDF文件
        
        Args:
            pdf_files (List[str]): PDF文件路径列表
            output_dir (str): 输出目录
            
        Returns:
            List[str]: 输出文件路径列表
        """
        
        results = []
        total_files = len(pdf_files)
        
        for i, pdf_path in enumerate(pdf_files):
            try:
                print(f"\n[{i+1}/{total_files}] 处理文件: {Path(pdf_path).name}")
                
                # 设置输出路径
                output_path = None
                if output_dir:
                    pdf_name = Path(pdf_path).stem
                    output_path = Path(output_dir) / f"{pdf_name}.pptx"
                
                # 转换文件
                result = self.convert_single_pdf(
                    pdf_path, 
                    str(output_path) if output_path else None,
                    progress_callback=lambda msg: print(f"  {msg}")
                )
                
                results.append(result)
                print(f"  ✅ 完成: {result}")
                
            except Exception as e:
                print(f"  ❌ 失败: {str(e)}")
                results.append(None)
        
        return results

def find_pdf_files(directory: str) -> List[str]:
    """查找目录中的所有PDF文件"""
    pdf_files = []
    for root, dirs, files in os.walk(directory):
        for file in files:
            if file.lower().endswith('.pdf'):
                pdf_files.append(os.path.join(root, file))
    return pdf_files

def main():
    """主函数"""
    parser = argparse.ArgumentParser(description='PDF to PowerPoint转换器')
    parser.add_argument('input', nargs='?', help='输入PDF文件或目录路径')
    parser.add_argument('-o', '--output', help='输出目录或文件路径')
    parser.add_argument('-d', '--dpi', type=int, default=200, help='图片分辨率 (默认: 200)')
    parser.add_argument('-p', '--poppler-path', help='Poppler工具路径 (Windows需要)')
    parser.add_argument('-b', '--batch', action='store_true', help='批量转换目录中的所有PDF文件')
    parser.add_argument('--version', action='version', version='PDF2PPT Converter 2.0')
    
    args = parser.parse_args()
    
    # 如果没有提供输入参数，使用默认文件
    if not args.input:
        script_dir = os.path.dirname(os.path.abspath(__file__))
        default_pdf = os.path.join(script_dir, "presentation.pdf")
        
        if os.path.exists(default_pdf):
            args.input = default_pdf
        else:
            print("错误: 请提供PDF文件路径或将PDF文件命名为 'presentation.pdf'")
            print("使用 --help 查看详细用法")
            sys.exit(1)
    
    # 创建转换器
    converter = PDFToPPTConverter(dpi=args.dpi, poppler_path=args.poppler_path)
    
    try:
        start_time = time.time()
        
        if args.batch or os.path.isdir(args.input):
            # 批量转换模式
            if not os.path.isdir(args.input):
                print("错误: 批量模式需要提供目录路径")
                sys.exit(1)
            
            print(f"正在搜索目录中的PDF文件: {args.input}")
            pdf_files = find_pdf_files(args.input)
            
            if not pdf_files:
                print("错误: 目录中没有找到PDF文件")
                sys.exit(1)
            
            print(f"找到 {len(pdf_files)} 个PDF文件")
            
            # 执行批量转换
            results = converter.convert_batch(pdf_files, args.output)
            
            # 统计结果
            successful = sum(1 for r in results if r is not None)
            failed = len(results) - successful
            
            print(f"\n批量转换完成:")
            print(f"  成功: {successful} 个文件")
            print(f"  失败: {failed} 个文件")
            
        else:
            # 单文件转换模式
            result = converter.convert_single_pdf(
                args.input, 
                args.output,
                progress_callback=lambda msg: print(msg)
            )
            print(f"\n转换成功! 输出文件: {result}")
        
        elapsed_time = time.time() - start_time
        print(f"总耗时: {elapsed_time:.2f} 秒")
        
    except Exception as e:
        print(f"\n转换失败: {str(e)}")
        print("\n故障排除提示:")
        print("1. 确保已安装poppler工具 (Windows)")
        print("2. 确保PDF文件路径正确")
        print("3. 确保有足够的磁盘空间和内存")
        print("4. 尝试降低DPI设置以减少内存使用")
        sys.exit(1)

if __name__ == "__main__":
    main()
