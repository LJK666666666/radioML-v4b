#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PDF to PowerPoint Converter - 智能尺寸版本
智能处理不同尺寸的PDF页面，为每个不同尺寸创建单独的PPT文件
"""

import os
import sys
import time
from pathlib import Path
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import tempfile
from collections import defaultdict

def analyze_pdf_pages(pdf_path, dpi=200):
    """
    分析PDF所有页面的尺寸
    
    Args:
        pdf_path (str): PDF文件路径
        dpi (int): 分析时使用的DPI
        
    Returns:
        list: 每页的尺寸信息 [(width_inch, height_inch, page_num), ...]
    """
    print("正在分析PDF页面尺寸...")
    
    # 先获取所有页面
    images = convert_from_path(pdf_path, dpi=dpi)
    page_info = []
    
    for i, image in enumerate(images):
        img_width, img_height = image.size
        width_inch = img_width / dpi
        height_inch = img_height / dpi
        page_info.append((round(width_inch, 2), round(height_inch, 2), i + 1))
        
        # 清理内存
        del image
    
    return page_info

def group_pages_by_size(page_info):
    """
    按照尺寸对页面进行分组
    
    Args:
        page_info: 页面信息列表
        
    Returns:
        dict: {(width, height): [page_numbers]}
    """
    size_groups = defaultdict(list)
    
    for width_inch, height_inch, page_num in page_info:
        size_key = (width_inch, height_inch)
        size_groups[size_key].append(page_num)
    
    return dict(size_groups)

def pdf_to_ppt_smart(pdf_path, output_dir=None, dpi=200):
    """
    智能PDF转PowerPoint转换器
    根据页面尺寸自动分组，为不同尺寸创建不同的PPT文件
    
    Args:
        pdf_path (str): PDF文件路径
        output_dir (str): 输出目录
        dpi (int): 图片分辨率，默认200
    """
    
    # 检查PDF文件是否存在
    if not os.path.exists(pdf_path):
        raise FileNotFoundError(f"PDF文件不存在: {pdf_path}")
    
    # 设置输出目录
    if output_dir is None:
        output_dir = os.path.dirname(pdf_path)
    
    pdf_name = Path(pdf_path).stem
    
    print(f"开始智能转换PDF文件: {pdf_path}")
    print(f"输出目录: {output_dir}")
    print(f"DPI设置: {dpi}")
    
    try:
        # 分析所有页面尺寸
        page_info = analyze_pdf_pages(pdf_path, dpi)
        print(f"PDF总页数: {len(page_info)}")
        
        # 按尺寸分组
        size_groups = group_pages_by_size(page_info)
        print(f"发现 {len(size_groups)} 种不同的页面尺寸:")
        
        for (width, height), pages in size_groups.items():
            print(f"  尺寸 {width}\" x {height}\": {len(pages)} 页 (页码: {pages})")
        
        output_files = []
        
        # 为每种尺寸创建PPT
        for size_idx, ((width_inch, height_inch), page_numbers) in enumerate(size_groups.items(), 1):
            
            if len(size_groups) == 1:
                # 如果只有一种尺寸，使用原文件名
                output_path = os.path.join(output_dir, f"{pdf_name}.pptx")
            else:
                # 如果有多种尺寸，添加尺寸后缀
                size_suffix = f"{width_inch}x{height_inch}in"
                output_path = os.path.join(output_dir, f"{pdf_name}_{size_suffix}.pptx")
            
            print(f"\n创建PPT文件 {size_idx}/{len(size_groups)}: {output_path}")
            print(f"幻灯片尺寸: {width_inch}\" x {height_inch}\"")
            print(f"包含页面: {page_numbers}")
            
            # 创建PowerPoint演示文稿
            prs = Presentation()
            
            # 设置幻灯片尺寸
            prs.slide_width = Inches(width_inch)
            prs.slide_height = Inches(height_inch)
            
            # 获取空白布局
            blank_slide_layout = prs.slide_layouts[6]
            
            # 转换指定页面
            for page_num in page_numbers:
                print(f"  处理第 {page_num} 页...")
                
                # 转换单页
                page_images = convert_from_path(
                    pdf_path, 
                    dpi=dpi, 
                    first_page=page_num, 
                    last_page=page_num
                )
                
                if not page_images:
                    print(f"    警告: 无法转换第 {page_num} 页")
                    continue
                
                image = page_images[0]
                
                # 添加幻灯片
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # 创建临时文件
                temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
                temp_file_path = temp_file.name
                temp_file.close()
                
                try:
                    # 保存图片
                    image.save(temp_file_path, 'PNG')
                    
                    # 图片占满整个幻灯片
                    slide.shapes.add_picture(
                        temp_file_path,
                        0,  # left
                        0,  # top
                        Inches(width_inch),   # width
                        Inches(height_inch)   # height
                    )
                    
                finally:
                    # 清理临时文件
                    for retry in range(3):
                        try:
                            if os.path.exists(temp_file_path):
                                os.unlink(temp_file_path)
                            break
                        except Exception as e:
                            if retry < 2:
                                time.sleep(0.1)
                                continue
                            else:
                                print(f"    警告: 无法删除临时文件: {e}")
                    
                    # 清理内存
                    del image
                    del page_images
            
            # 保存PPT文件
            prs.save(output_path)
            output_files.append(output_path)
            print(f"  ✅ 保存完成: {output_path}")
        
        print(f"\n🎉 转换完成!")
        print(f"生成的PPT文件:")
        for output_file in output_files:
            print(f"  - {output_file}")
        
        return output_files
        
    except Exception as e:
        print(f"转换过程中出现错误: {str(e)}")
        raise

def main():
    """主函数"""
    # 获取当前脚本所在目录
    script_dir = os.path.dirname(os.path.abspath(__file__))
    
    # 默认PDF文件路径
    default_pdf_path = os.path.join(script_dir, "presentation.pdf")
    
    # 检查命令行参数
    if len(sys.argv) > 1:
        pdf_path = sys.argv[1]
        dpi = int(sys.argv[2]) if len(sys.argv) > 2 else 200
        output_dir = sys.argv[3] if len(sys.argv) > 3 else None
    else:
        pdf_path = default_pdf_path
        dpi = 200
        output_dir = None
    
    # 检查PDF文件是否存在
    if not os.path.exists(pdf_path):
        print(f"错误: PDF文件不存在: {pdf_path}")
        print("用法: python pdf_to_ppt_smart.py [PDF文件路径] [DPI(可选)] [输出目录(可选)]")
        print(f"或者将PDF文件命名为 'presentation.pdf' 并放在脚本同目录下")
        sys.exit(1)
    
    try:
        start_time = time.time()
        
        # 执行智能转换
        result_files = pdf_to_ppt_smart(pdf_path, output_dir, dpi)
        
        end_time = time.time()
        elapsed_time = end_time - start_time
        
        print(f"\n✨ 转换成功完成！")
        print(f"⏱️  耗时: {elapsed_time:.2f} 秒")
        print(f"📁 生成了 {len(result_files)} 个PPT文件")
        
    except Exception as e:
        print(f"\n❌ 转换失败: {str(e)}")
        print("\n🔧 故障排除提示:")
        print("1. 确保PDF文件未被其他程序占用")
        print("2. 确保有足够的磁盘空间")
        print("3. 尝试降低DPI设置")
        print("4. 确保已正确安装所有依赖包")
        sys.exit(1)

if __name__ == "__main__":
    main()
